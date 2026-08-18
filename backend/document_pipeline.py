from __future__ import annotations

import hashlib
import io
import math
import mimetypes
import os
import re
import unicodedata
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable


EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "BAAI/bge-small-en-v1.5")
# Cross-encoder used to rescore retrieval candidates. It reads the query and the
# passage together, so unlike a bi-encoder cosine score its output is comparable
# across queries and can be thresholded directly.
RERANK_MODEL = os.getenv("RERANK_MODEL", "Xenova/ms-marco-MiniLM-L-6-v2")
# Matched CLIP pair sharing one embedding space, used to understand images that
# contain no readable text. Without this, a photograph has no searchable content.
# Measured against real photos, jina-clip-v1 was more confident but less accurate at
# object naming than OpenAI CLIP ViT-B/32: it called a backpack a chair at 0.78, two
# computer accessories tables at 0.99, and a medical report a sofa. Ranking was wrong, not
# just calibration, so a bigger model here would have cost a schema migration and lost
# accuracy. Keeping ViT-B/32.
IMAGE_VISION_MODEL = os.getenv("IMAGE_VISION_MODEL", "Qdrant/clip-ViT-B-32-vision")
IMAGE_TEXT_MODEL = os.getenv("IMAGE_TEXT_MODEL", "Qdrant/clip-ViT-B-32-text")
# CLIP's trained logit scale. jina-clip produces a different similarity distribution to
# OpenAI CLIP, so this is configurable rather than hard-coded at 100.
IMAGE_LOGIT_SCALE = float(os.getenv("IMAGE_LOGIT_SCALE", "100"))
# Zero-shot labelling is only trusted above this probability, and it must beat the
# runner-up by this margin, otherwise the image is recorded as a generic photo.
IMAGE_LABEL_THRESHOLD = 0.18
IMAGE_LABEL_MARGIN = 0.04
MAX_FILE_BYTES = int(os.getenv("MAX_INGESTION_FILE_BYTES", str(25 * 1024 * 1024)))
MAX_PDF_PAGES = int(os.getenv("MAX_PDF_PAGES", "80"))
MAX_CHUNKS = int(os.getenv("MAX_DOCUMENT_CHUNKS", "200"))
OCR_MIN_CHARS_PER_PAGE = 80

# Classification thresholds apply to the blended rule/semantic/lexical score in
# classify_document. Raw embedding similarity between a long document and a short
# folder label clusters in a narrow band, so it is normalised by spread and only
# used to break ties. Deterministic rule evidence carries the routing decision.
AUTO_CLASSIFICATION_THRESHOLD = 0.50
AUTO_CLASSIFICATION_MARGIN = 0.10
# Alternatives below this score are noise and are never offered to the user.
ALTERNATIVE_SCORE_FLOOR = 0.32
# Without rule evidence a folder must be a strong similarity match to be suggested.
STRONG_SIMILARITY_SCORE = 0.45
# Embedding scores are ignored entirely when every folder scores about the same.
MIN_SEMANTIC_SPREAD = 0.02


class DocumentPipelineError(Exception):
    """Base error for local document processing."""


class FileTooLargeError(DocumentPipelineError):
    pass


@dataclass
class TextChunk:
    index: int
    content: str
    page_number: int | None = None


@dataclass
class ExtractedDocument:
    text: str
    chunks: list[TextChunk]
    extraction_method: str
    extraction_status: str
    page_count: int | None = None
    extraction_error: str | None = None
    # Heading detected from OCR layout, when the source was an image.
    heading: str | None = None


@dataclass
class DocumentMetadata:
    title: str
    summary: str
    document_type: str
    keywords: list[str]
    entities: dict[str, Any]
    language: str = "unknown"


@dataclass
class FolderCandidate:
    folder_id: str
    label: str
    subject: str
    unit: str | None
    score: float = 0.0


@dataclass
class ClassificationResult:
    target_folder_id: str
    target_label: str
    subject: str
    unit: str | None
    confidence: float
    status: str
    alternatives: list[FolderCandidate] = field(default_factory=list)


_STOP_WORDS = {
    "about", "after", "again", "also", "and", "are", "been", "before", "being", "between",
    "but", "can", "chapter", "class", "course", "document", "each", "file", "for", "from",
    "have", "into", "its", "more", "notes", "page", "part", "that", "the", "their", "these",
    "this", "those", "through", "unit", "using", "was", "were", "what", "when", "where",
    "which", "will", "with", "your",
}

_MIME_BY_EXTENSION = {
    ".pdf": "application/pdf",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".gif": "image/gif",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
    ".json": "application/json",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


_embedder = None
_ocr_engine = None
_reranker = None
_image_encoder = None
_image_text_encoder = None
_prompt_vectors: dict[str, list[float]] = {}
# Metadata building and indexing both need the same image analysis, so the most
# recent result is memoised to avoid running vision inference twice per file.
_last_image_analysis: tuple[str, "ImageDescription | None"] | None = None


@dataclass
class ImageDescription:
    """What an image appears to show, plus its CLIP vector for visual search."""
    label: str
    document_type: str
    caption: str
    keywords: list[str]
    confidence: float
    vector: list[float] | None = None


# Zero-shot concepts for CLIP. Each entry is a prompt, the signal it maps to, the
# document type to record, a human caption, and keywords for full-text search.
_IMAGE_CONCEPTS: tuple[tuple[str, str, str, str, tuple[str, ...]], ...] = (
    (
        "a portrait photograph of a single person's face",
        "person_photo",
        "person_photo",
        "Portrait photo of a person",
        ("photo", "portrait", "person", "face", "people"),
    ),
    (
        "a photograph of a group of people together at an event",
        "event_photo",
        "event_photo",
        "Photo of a group of people",
        ("photo", "group", "people", "event", "gathering"),
    ),
    (
        "a screenshot of a mobile app or computer software interface",
        "screenshot",
        "screenshot",
        "Screenshot of an app interface",
        ("screenshot", "app", "interface", "screen"),
    ),
    (
        "a scanned page of a printed text document",
        "study_notes",
        "document_scan",
        "Scanned printed document page",
        ("document", "scan", "page", "printed"),
    ),
    (
        "a photograph of handwritten notes on paper",
        "study_notes",
        "handwritten_notes",
        "Photo of handwritten notes",
        ("handwritten", "notes", "paper", "writing"),
    ),
    (
        "a photograph of a whiteboard or blackboard with writing on it",
        "study_notes",
        "whiteboard_photo",
        "Photo of a whiteboard",
        ("whiteboard", "board", "notes", "lecture"),
    ),
    (
        "a photograph of a government identity card or licence",
        "identity_other",
        "identity_card",
        "Photo of an identity card",
        ("identity", "card", "government", "id"),
    ),
    (
        "a photograph of a paper receipt or a printed bill",
        "invoice",
        "receipt",
        "Photo of a receipt or bill",
        ("receipt", "bill", "invoice", "payment"),
    ),
    (
        "a photograph of an outdoor landscape, building or scenery",
        "scenery_photo",
        "scenery_photo",
        "Photograph of scenery",
        ("photo", "scenery", "landscape", "outdoor", "place"),
    ),
    (
        "a photograph of a meal, food or a drink",
        "food_photo",
        "food_photo",
        "Photo of food",
        ("photo", "food", "meal", "drink"),
    ),
    (
        "a photograph of a car, bike or vehicle",
        "vehicle_photo",
        "vehicle_photo",
        "Photo of a vehicle",
        ("photo", "vehicle", "car", "bike"),
    ),
    (
        "a product photograph of an object on a plain background",
        "product_photo",
        "product_photo",
        "Product photo of an object",
        ("photo", "product", "object", "item"),
    ),
    (
        "a medical prescription, x-ray scan or lab report",
        "medical_report",
        "medical_report",
        "Medical report or prescription",
        ("medical", "health", "report", "prescription"),
    ),
    (
        "a chart, graph or data visualisation",
        "chart_image",
        "chart",
        "Chart or data visualisation",
        ("chart", "graph", "data", "visualisation"),
    ),
    (
        "a certificate, diploma or award document",
        "certificate",
        "certificate",
        "Certificate or award document",
        ("certificate", "award", "diploma"),
    ),
)

# Naming the thing in a photo. The coarse concept list above cannot distinguish a chair
# from a backpack, so every object became "Product photo of an object" with no searchable
# noun, and text search could never find it. Worse, a missing option forced images into
# the nearest wrong bucket: a backpack was labelled a vehicle.
#
# Asking which of many nouns fits one image is far better conditioned than asking which
# of many images fits one query, so this is where the vision model is actually reliable.
# Each entry is a prompt, the display name, the document type it implies, and the
# keywords that make it findable, including synonyms a user might search instead.
_OBJECT_NOUNS: tuple[tuple[str, str, str, tuple[str, ...]], ...] = (
    ("a photo of a chair or an office chair", "chair", "product_photo", ("chair", "seat", "furniture")),
    ("a photo of a table or a desk", "table", "product_photo", ("table", "desk", "furniture")),
    ("a photo of a bed", "bed", "product_photo", ("bed", "mattress", "furniture")),
    ("a photo of a sofa or couch", "sofa", "product_photo", ("sofa", "couch", "furniture")),
    ("a photo of a door", "door", "product_photo", ("door", "doorway", "entrance")),
    ("a photo of a window", "window", "product_photo", ("window", "glass")),
    ("a photo of a notebook, diary or notepad", "notebook", "product_photo", ("notebook", "diary", "notepad", "journal")),
    ("a photo of a book", "book", "product_photo", ("book", "textbook", "novel")),
    ("a photo of a backpack, bag or luggage", "bag", "product_photo", ("bag", "backpack", "rucksack", "luggage", "handbag")),
    ("a photo of earphones, earbuds or headphones", "earphones", "product_photo", ("earphones", "earbuds", "headphones", "headset", "audio")),
    ("a photo of a laptop or computer", "laptop", "product_photo", ("laptop", "computer", "notebook")),
    ("a photo of a computer monitor or television screen", "monitor", "product_photo", ("monitor", "screen", "display", "television")),
    ("a photo of a keyboard or a computer mouse", "computer accessory", "product_photo", ("keyboard", "mouse", "accessory", "peripheral")),
    ("a photo of a mobile phone", "mobile phone", "product_photo", ("phone", "mobile", "smartphone")),
    ("a photo of a charger, cable or wires", "cable", "product_photo", ("cable", "charger", "wire", "wires", "cord")),
    ("a photo of a wristwatch", "watch", "product_photo", ("watch", "wristwatch", "timepiece")),
    ("a photo of spectacles or eyeglasses on a surface", "spectacles", "product_photo", ("spectacles", "glasses", "eyeglasses", "specs")),
    ("a photo of shoes or footwear", "shoes", "product_photo", ("shoes", "footwear", "sneakers")),
    ("a photo of clothing or a shirt on its own", "clothing", "product_photo", ("clothing", "clothes", "shirt", "apparel")),
    ("a photo of a water bottle or flask", "bottle", "product_photo", ("bottle", "flask", "water")),
    ("a photo of a cup of coffee or tea", "cup", "food_photo", ("cup", "coffee", "tea", "mug", "drink")),
    ("a photo of a plate of cooked food or a meal", "food", "food_photo", ("food", "meal", "plate", "dish")),
    ("a photo of fruit or vegetables", "fruit", "food_photo", ("fruit", "vegetables", "produce", "food")),
    ("a photo of a car", "car", "vehicle_photo", ("car", "vehicle", "automobile")),
    ("a photo of a motorcycle or scooter", "motorcycle", "vehicle_photo", ("motorcycle", "bike", "scooter", "vehicle")),
    ("a photo of a bicycle", "bicycle", "vehicle_photo", ("bicycle", "cycle", "bike")),
    ("a photo of a houseplant or flowers", "plant", "product_photo", ("plant", "flowers", "greenery")),
    ("a photo of a wallet or purse", "wallet", "product_photo", ("wallet", "purse", "money")),
    ("a photo of keys or a keychain", "keys", "product_photo", ("keys", "keychain", "key")),
    ("a photo of a fan, lamp or light fixture", "light or fan", "product_photo", ("lamp", "light", "fan", "fixture")),
    ("a photo of a bottle of medicine or pills", "medicine", "product_photo", ("medicine", "pills", "tablets", "medical")),
    ("a photo of a pen or pencil", "pen", "product_photo", ("pen", "pencil", "stationery")),
    ("a photo of a room interior", "room", "scenery_photo", ("room", "interior", "indoor", "space")),
    ("a photo of a building from outside", "building", "scenery_photo", ("building", "exterior", "architecture")),
    ("a photo of trees, sky or natural landscape", "landscape", "scenery_photo", ("landscape", "nature", "outdoor", "scenery")),
)

# Coarse labels vague enough that naming the object is an improvement. People and
# documents are excluded: they have their own, better handling.
_OBJECT_PROBE_LABELS = {"product_photo", "photo", "food_photo", "vehicle_photo", "scenery_photo"}

# With this many prompts, a uniform guess is under 0.03, so these bars are well above
# chance while still admitting the honest uncertainty seen in testing.
OBJECT_LABEL_THRESHOLD = 0.15
OBJECT_LABEL_MARGIN = 0.04

# Runner-up synonyms are deliberately NOT indexed. Measured on real photos, the runner-up
# was either irrelevant (ratio under 0.15) or actively wrong: a backpack scored "chair" at
# 0.55 of the winner, which made a search for "chair" confidently return the bag. Missing
# a file is recoverable; confidently returning the wrong one is not.


# Attribute probing. Fixed concept labels alone give every photo of the same kind an
# identical title, which makes two portraits impossible to tell apart. Each group below
# is a set of mutually exclusive options scored against one another, so the winner adds
# a distinguishing detail. Options with an empty fragment are valid outcomes that simply
# contribute nothing to the caption.
_ATTRIBUTE_GROUPS: dict[str, tuple[tuple[str, str, tuple[str, ...]], ...]] = {
    "subject_count": (
        ("a photo of exactly one person alone", "", ()),
        ("a photo of two people together", "two people", ("two", "pair", "couple")),
        ("a photo of a large group of many people", "a group of people", ("group", "many", "crowd")),
    ),
    "age_gender": (
        ("a photo of a young man in his twenties", "a young man", ("young", "man", "male", "guy")),
        ("a photo of a middle aged or elderly man", "an older man", ("older", "elderly", "man", "male")),
        ("a photo of a young woman in her twenties", "a young woman", ("young", "woman", "female", "girl")),
        ("a photo of a middle aged or elderly woman", "an older woman", ("older", "elderly", "woman", "female")),
        ("a photo of a young child", "a child", ("child", "kid", "young")),
    ),
    "eyewear": (
        ("a person wearing eyeglasses on their face", "wearing glasses", ("glasses", "eyeglasses", "spectacles", "specs")),
        ("a person with no eyeglasses", "", ()),
    ),
    "facial_hair": (
        ("a man with a beard or moustache", "with a beard", ("beard", "moustache", "facial")),
        ("a clean shaven face with no beard", "", ()),
    ),
    "shot_type": (
        ("a close up selfie taken at arms length", "selfie", ("selfie", "closeup", "close")),
        # Deliberately avoids words like "passport" or "identity": this is a photo of a
        # person in a formal pose, and those terms would be misread as evidence that the
        # file is an actual identity document and route it into an ID folder.
        ("a formal identity document style headshot on a plain background", "formal headshot", ("headshot", "formal", "official")),
        ("a casual candid photo of a person", "", ()),
    ),
    "clothing": (
        ("a person wearing a white shirt", "in a white shirt", ("white", "shirt")),
        ("a person wearing a black shirt", "in a black shirt", ("black", "shirt")),
        ("a person wearing a brown shirt", "in a brown shirt", ("brown", "shirt")),
        ("a person wearing a blue shirt", "in a blue shirt", ("blue", "shirt")),
        ("a person wearing a grey shirt", "in a grey shirt", ("grey", "shirt")),
        ("a person wearing a red shirt", "in a red shirt", ("red", "shirt")),
        ("a person wearing a green shirt", "in a green shirt", ("green", "shirt")),
        ("a person wearing a yellow or orange shirt", "in a yellow shirt", ("yellow", "orange", "shirt")),
    ),
    "setting": (
        ("a photo taken indoors inside a room", "indoors", ("indoor", "inside", "room")),
        ("a photo taken outdoors outside in daylight", "outdoors", ("outdoor", "outside")),
        ("a photo against a plain studio backdrop", "", ()),
    ),
}

# Which attribute groups make sense for which primary concept. Probing clothing on a
# dashboard screenshot would only produce nonsense.
_ATTRIBUTE_GROUPS_BY_LABEL: dict[str, tuple[str, ...]] = {
    "person_photo": ("subject_count", "age_gender", "eyewear", "facial_hair", "shot_type", "clothing", "setting"),
    "event_photo": ("subject_count", "setting"),
    "scenery_photo": ("setting",),
    "photo": ("subject_count", "setting"),
    # Where an object was photographed is often how a user remembers it.
    "product_photo": ("setting",),
    "food_photo": ("setting",),
    "vehicle_photo": ("setting",),
}

# An attribute is only recorded when it clearly wins its group, so captions stay honest.
ATTRIBUTE_THRESHOLD = 0.45
ATTRIBUTE_MARGIN = 0.10


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def detect_mime_type(data: bytes, claimed_mime: str | None, filename: str) -> str:
    """Prefer signatures over untrusted WhatsApp metadata and extensions."""
    if data.startswith(b"%PDF"):
        return "application/pdf"
    if data.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if data.startswith((b"GIF87a", b"GIF89a")):
        return "image/gif"
    if data.startswith(b"RIFF") and data[8:12] == b"WEBP":
        return "image/webp"

    extension = Path(filename or "").suffix.lower()
    if data.startswith(b"PK\x03\x04") and extension in {".docx", ".pptx", ".xlsx"}:
        return _MIME_BY_EXTENSION[extension]

    if extension in _MIME_BY_EXTENSION:
        return _MIME_BY_EXTENSION[extension]
    if claimed_mime and claimed_mime != "application/octet-stream":
        return claimed_mime.split(";", 1)[0].strip().lower()
    return mimetypes.guess_type(filename)[0] or "application/octet-stream"


def safe_filename(original_filename: str, mime_type: str) -> str:
    name = Path(original_filename or "document").name
    stem = unicodedata.normalize("NFKC", Path(name).stem)
    stem = re.sub(r"[^\w .()\-]+", "_", stem, flags=re.UNICODE)
    stem = re.sub(r"\s+", " ", stem).strip(" ._-") or "document"
    extension = Path(name).suffix.lower()
    if not extension:
        extension = mimetypes.guess_extension(mime_type) or ""
    return f"{stem[:120]}{extension[:12]}"


def _normalize_text(text: str) -> str:
    text = unicodedata.normalize("NFKC", text or "")
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def split_glued_words(text: str) -> str:
    """Insert the spaces OCR dropped between words.

    Region-based OCR frequently returns a heading as one run of characters, so
    "MemoryHierarchyDesign" arrives as a single token. That breaks folder matching and
    makes the text unsearchable: a query for "memory hierarchy" cannot match it, while
    "memoryhierarchydesign" scores 0.999. Capitalisation survives OCR, so case
    boundaries recover most of the missing spaces.
    """
    if not text:
        return text
    # "designCPU" -> "design CPU", and "hierarchyDesign" -> "hierarchy Design".
    repaired = re.sub(r"(?<=[a-z0-9])(?=[A-Z])", " ", text)
    # "CPURegisters" -> "CPU Registers", keeping acronyms intact.
    repaired = re.sub(r"(?<=[A-Z])(?=[A-Z][a-z])", " ", repaired)
    return re.sub(r"[ \t]+", " ", repaired)


def _get_ocr_engine():
    global _ocr_engine
    if _ocr_engine is None:
        from rapidocr_onnxruntime import RapidOCR

        _ocr_engine = RapidOCR()
    return _ocr_engine


def _ocr_image_bytes(image_bytes: bytes) -> str:
    import numpy as np
    from PIL import Image

    text, _ = _ocr_image_regions(image_bytes)
    return text


def _ocr_image_regions(image_bytes: bytes) -> tuple[str, str | None]:
    """OCR an image and also pick out its heading.

    RapidOCR returns a bounding box per text region, and box height stands in for font
    size. The largest text in the upper part of an image is almost always the heading,
    which is far more reliable than taking the first line in reading order.
    """
    import numpy as np
    from PIL import Image

    image = np.asarray(Image.open(io.BytesIO(image_bytes)).convert("RGB"))
    result, _ = _get_ocr_engine()(image)
    if not result:
        return "", None

    image_height = image.shape[0] or 1
    lines: list[str] = []
    regions: list[tuple[float, float, str]] = []

    for line in result:
        if len(line) < 2:
            continue
        content = str(line[1]).strip()
        if not content:
            continue
        lines.append(content)

        box = line[0]
        try:
            ys = [float(point[1]) for point in box]
            xs = [float(point[0]) for point in box]
            height = max(ys) - min(ys)
            width = max(xs) - min(xs)
            top = min(ys) / image_height
        except (TypeError, IndexError, ValueError):
            continue
        regions.append((height, width, top, content))

    text = _normalize_text("\n".join(lines))
    return text, _heading_from_regions(regions)


def _heading_from_regions(regions: list[tuple[float, float, float, str]]) -> str | None:
    """Return a heading only when the layout makes it unambiguous.

    Deliberately conservative, because looser versions did more harm than good. Rotated
    labels in diagrams produce tall narrow boxes that look like huge text, so a candidate
    must read horizontally, sit near the top, and be the largest text present before it
    overrides simple reading order.
    """
    horizontal = [
        (height, top, content)
        for height, width, top, content in regions
        # A line of horizontal text is always wider than it is tall. Vertical axis labels
        # are not, and one of those was previously chosen as a document title.
        if height > 0 and width / height >= 1.5
    ]
    if not horizontal:
        return None

    tallest = max(height for height, _, _ in horizontal)
    if tallest <= 0:
        return None

    for height, top, content in sorted(horizontal, key=lambda item: item[1]):
        words = re.findall(r"[A-Za-z]{2,}", content)
        # Headings are short. A long line of body text is not a title even if the glyphs
        # happen to be large.
        if not words or len(content) > 90:
            continue
        if top <= 0.30 and height / tallest >= 0.90:
            return content
    return None


def _extract_pdf(data: bytes) -> tuple[list[tuple[int | None, str]], str, int]:
    import pymupdf

    pages: list[tuple[int | None, str]] = []
    used_ocr = False
    with pymupdf.open(stream=data, filetype="pdf") as document:
        page_count = min(document.page_count, MAX_PDF_PAGES)
        for page_index in range(page_count):
            page = document.load_page(page_index)
            text = _normalize_text(page.get_text("text"))
            if len(re.sub(r"\W", "", text)) < OCR_MIN_CHARS_PER_PAGE:
                pixmap = page.get_pixmap(matrix=pymupdf.Matrix(2, 2), alpha=False)
                ocr_text = _ocr_image_bytes(pixmap.tobytes("png"))
                if len(ocr_text) > len(text):
                    text = ocr_text
                    used_ocr = True
            if text:
                pages.append((page_index + 1, text))
    method = "pdf_text+ocr" if used_ocr else "pdf_text"
    return pages, method, page_count


def _extract_docx(data: bytes) -> list[tuple[int | None, str]]:
    from docx import Document

    document = Document(io.BytesIO(data))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                parts.append(" | ".join(values))
    return [(None, _normalize_text("\n".join(parts)))]


def _extract_pptx(data: bytes) -> list[tuple[int | None, str]]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    pages: list[tuple[int | None, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        text = _normalize_text("\n".join(texts))
        if text:
            pages.append((index, text))
    return pages


def _extract_xlsx(data: bytes) -> list[tuple[int | None, str]]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    pages: list[tuple[int | None, str]] = []
    for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        text = _normalize_text(f"{sheet.title}\n" + "\n".join(rows))
        if text:
            pages.append((sheet_index, text))
    return pages


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-16", "latin-1"):
        try:
            return _normalize_text(data.decode(encoding))
        except UnicodeDecodeError:
            continue
    return ""


def _split_section(text: str, target_size: int = 1000, overlap: int = 150) -> list[str]:
    text = _normalize_text(text)
    if not text:
        return []
    if len(text) <= target_size:
        return [text]

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + target_size, len(text))
        if end < len(text):
            boundary = max(text.rfind("\n", start, end), text.rfind(". ", start, end))
            if boundary > start + target_size // 2:
                end = boundary + 1
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = max(end - overlap, start + 1)
    return chunks


def extract_document(data: bytes, mime_type: str, filename: str) -> ExtractedDocument:
    heading: str | None = None
    if len(data) > MAX_FILE_BYTES:
        raise FileTooLargeError(f"File exceeds the {MAX_FILE_BYTES // (1024 * 1024)} MB limit")

    pages: list[tuple[int | None, str]] = []
    method = "metadata_only"
    page_count: int | None = None

    if mime_type == "application/pdf":
        pages, method, page_count = _extract_pdf(data)
    elif mime_type.startswith("image/"):
        raw_text, raw_heading = _ocr_image_regions(data)
        # Only OCR output needs this. Extracted PDF and Office text already has real
        # spaces, where splitting would wrongly break names like "JavaScript".
        text = split_glued_words(raw_text)
        heading = split_glued_words(raw_heading).strip() if raw_heading else None
        pages = [(1, text)] if text else []
        method, page_count = "image_ocr", 1
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        pages, method = _extract_docx(data), "docx_text"
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        pages, method = _extract_pptx(data), "pptx_text"
        page_count = len(pages)
    elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        pages, method = _extract_xlsx(data), "xlsx_text"
        page_count = len(pages)
    elif mime_type.startswith("text/") or mime_type in {"application/json", "application/xml"}:
        text = _decode_text(data)
        pages = [(None, text)] if text else []
        method = "plain_text"

    chunks: list[TextChunk] = []
    full_text_parts = []
    for page_number, page_text in pages:
        if not page_text:
            continue
        full_text_parts.append(page_text)
        for content in _split_section(page_text):
            if len(chunks) >= MAX_CHUNKS:
                break
            chunks.append(TextChunk(index=len(chunks), content=content, page_number=page_number))
        if len(chunks) >= MAX_CHUNKS:
            break

    full_text = _normalize_text("\n\n".join(full_text_parts))
    extraction_status = "complete" if full_text else ("unsupported" if method == "metadata_only" else "no_text")
    return ExtractedDocument(
        text=full_text,
        chunks=chunks,
        extraction_method=method,
        extraction_status=extraction_status,
        page_count=page_count,
        heading=heading,
    )


def _tokens(text: str) -> list[str]:
    return [
        token.lower()
        for token in re.findall(r"[A-Za-z][A-Za-z0-9+#.\-]{2,}", text or "")
        if token.lower() not in _STOP_WORDS
    ]


def _extract_keywords(text: str, limit: int = 15) -> list[str]:
    counts = Counter(_tokens(text))
    return [token for token, _ in counts.most_common(limit)]


_PAN_PATTERN = re.compile(r"\b[A-Z]{5}[0-9]{4}[A-Z]\b")
_AADHAAR_PATTERN = re.compile(r"\b\d{4}\s?\d{4}\s?\d{4}\b")

# A screenshot is defined by how it was captured, not by what it says, so it needs
# interface vocabulary rather than topical matching.
_SCREENSHOT_STRONG_TERMS = (
    "dashboard", "screenshot", "screen shot", "settings", "notifications",
    "log in", "login", "sign in", "sign up", "inbox", "my account",
)
_SCREENSHOT_WEAK_TERMS = (
    "search", "menu", "profile", "home", "active", "pending", "total",
    "view all", "see all", "share", "download", "upload", "filter", "sort",
    "app", "tap", "swipe", "online", "offline", "balance", "overview",
)


# Source code captured as an image is still a screenshot, but it shares no vocabulary
# with an app interface, so it needs its own indicators.
_CODE_TERMS = (
    "#include", "int main", "public static", "def ", "function ", "console.log",
    "vector<", "std::", "printf", "system.out", "getline", "return 0", "class ",
    "import ", "const ", "let ", "var ", "null", "void ", "elif", "lambda",
)


def _looks_like_source_code(corpus: str) -> bool:
    hits = sum(1 for term in _CODE_TERMS if term in corpus)
    symbols = sum(corpus.count(character) for character in ("{", "}", ";", "()"))
    return hits >= 2 or (hits >= 1 and symbols >= 6)


def _looks_like_screenshot(corpus: str) -> bool:
    if any(term in corpus for term in _SCREENSHOT_STRONG_TERMS):
        return True
    if _looks_like_source_code(corpus):
        return True
    return sum(1 for term in _SCREENSHOT_WEAK_TERMS if term in corpus) >= 3


def _identity_subtype(corpus: str, raw: str) -> str | None:
    """Return a precise identity signal, or None when the document is not an ID."""
    if "aadhaar" in corpus or "aadhar" in corpus or "unique identification authority" in corpus:
        return "aadhaar"
    if "permanent account number" in corpus or "pan card" in corpus:
        return "pan"
    if _PAN_PATTERN.search(raw) and ("income tax" in corpus or "permanent account" in corpus):
        return "pan"
    if "passport" in corpus:
        return "passport"
    if "driving licence" in corpus or "driving license" in corpus:
        return "driving_licence"
    if "voter id" in corpus or "election commission" in corpus:
        return "voter_id"
    # A bare 12-digit group is only an Aadhaar number alongside ID context.
    if _AADHAAR_PATTERN.search(raw) and any(
        term in corpus for term in ("government of india", "date of birth", "dob", "yob", "father")
    ):
        return "aadhaar"
    return None


def _guess_document_type(text: str, mime_type: str, filename: str) -> str:
    corpus = f"{filename} {text[:5000]}".lower()
    raw = f"{filename} {text[:5000]}"

    if _identity_subtype(corpus, raw):
        return "identity_card"

    patterns = {
        "receipt": ("receipt", "subtotal", "invoice", "amount due", "tax invoice"),
        "bank_statement": ("account statement", "bank statement", "closing balance", "ifsc"),
        "syllabus": ("syllabus", "course outcomes", "course objectives", "module i", "unit i"),
        "assignment": ("assignment", "submit by", "due date", "homework"),
        "question_paper": ("question paper", "maximum marks", "time allowed", "answer all"),
        "certificate": ("certificate", "hereby certify", "awarded to"),
        "resume": ("curriculum vitae", "work experience", "skills"),
        "medical_report": ("prescription", "diagnosis", "patient name", "dosage"),
        "study_notes": ("lecture notes", "chapter", "module"),
    }
    for document_type, signals in patterns.items():
        if any(signal in corpus for signal in signals):
            return document_type

    if mime_type.startswith("image/"):
        return "screenshot" if _looks_like_screenshot(corpus) else "image"
    if mime_type == "application/pdf":
        return "pdf_document"
    return "document"


def _clean_heading_text(value: str | None) -> str | None:
    """Accept a layout heading only if it reads like a real title."""
    candidate = re.sub(r"\s+", " ", value or "").strip(" .,:;|-–_")
    if len(candidate) < 4 or len(candidate) > 90:
        return None
    words = re.findall(r"[A-Za-z]{2,}", candidate)
    if not words:
        return None
    # Reject rows that are mostly digits, such as a clock or a table of numbers.
    letters = sum(character.isalpha() for character in candidate)
    if letters < len(re.sub(r"\s", "", candidate)) * 0.5:
        return None
    return candidate[:120]


def _image_text_is_sparse(text: str) -> bool:
    """True when an image's OCR output is too thin to describe the image."""
    stripped = (text or "").strip()
    words = re.findall(r"[A-Za-z]{3,}", stripped)
    return len(stripped) < 60 or len(words) < 4


def _title_from_text(text: str, filename: str) -> str:
    for line in text.splitlines()[:20]:
        candidate = re.sub(r"\s+", " ", line).strip(" -_:|")
        if 4 <= len(candidate) <= 120 and len(candidate.split()) <= 16:
            if not re.fullmatch(r"[\d\W]+", candidate):
                return candidate
    stem = re.sub(r"[_\-]+", " ", Path(filename).stem).strip()
    return stem[:120] or "Document"


def build_metadata(
    extracted: ExtractedDocument,
    filename: str,
    mime_type: str,
    image_data: bytes | None = None,
) -> DocumentMetadata:
    text = extracted.text

    # A photograph yields little or no OCR text, so fall back to visual labelling.
    # Sparse text is treated the same as none, because photos frequently produce a few
    # stray characters from clothing or background signage, and trusting those would
    # skip visual understanding and leave the file effectively unsearchable.
    if image_data and mime_type.startswith("image/") and _image_text_is_sparse(text):
        description = describe_image(image_data)
        if description:
            fragment = re.sub(r"\s+", " ", text).strip()[:200]
            summary = f"{description.caption}. Recognised from image content."
            if fragment:
                summary = f"{summary} Text found in image: {fragment}"
            return DocumentMetadata(
                title=description.caption,
                summary=summary[:600],
                document_type=description.document_type,
                keywords=description.keywords + [
                    keyword for keyword in _extract_keywords(text, limit=8)
                    if keyword not in description.keywords
                ],
                entities={"emails": [], "dates": []},
                language="unknown",
            )

    # A heading found by layout beats the first line in reading order, which is how
    # titles like "ntelli Paat" and "Team Name Team S Po C" were produced.
    title = _clean_heading_text(extracted.heading) or _title_from_text(text, filename)
    meaningful_lines = [line.strip() for line in text.splitlines() if len(line.strip()) > 30]
    summary_source = " ".join(meaningful_lines[:4]) or text[:600]
    summary = re.sub(r"\s+", " ", summary_source).strip()[:600]
    keywords = _extract_keywords(f"{title}\n{text}")

    entities = {
        "emails": sorted(set(re.findall(r"[\w.+-]+@[\w.-]+\.[A-Za-z]{2,}", text)))[:10],
        "dates": sorted(set(re.findall(
            r"\b(?:\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|\d{1,2}\s+(?:Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|Jul(?:y)?|Aug(?:ust)?|Sep(?:tember)?|Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)\s+\d{2,4})\b",
            text,
            flags=re.IGNORECASE,
        )))[:20],
    }
    return DocumentMetadata(
        title=title,
        summary=summary,
        document_type=_guess_document_type(text, mime_type, filename),
        keywords=keywords,
        entities=entities,
        language="en" if text and sum(ch.isascii() for ch in text) / max(len(text), 1) > 0.85 else "unknown",
    )


def _get_embedder():
    global _embedder
    if _embedder is None:
        from fastembed import TextEmbedding

        _embedder = TextEmbedding(model_name=EMBEDDING_MODEL)
    return _embedder


def embed_texts(texts: Iterable[str], *, query: bool = False) -> list[list[float] | None]:
    values = [re.sub(r"\s+", " ", text or " ").strip()[:12000] or "empty document" for text in texts]
    if not values:
        return []
    try:
        model = _get_embedder()
        generator = model.query_embed(values) if query and hasattr(model, "query_embed") else model.embed(values)
        return [[float(value) for value in vector] for vector in generator]
    except Exception as exc:
        print(f"Local embeddings unavailable; continuing with full-text search: {exc}")
        return [None for _ in values]


def embed_query(query: str) -> list[float] | None:
    return embed_texts([query], query=True)[0]


def _get_image_encoders():
    global _image_encoder, _image_text_encoder
    if _image_encoder is None or _image_text_encoder is None:
        from fastembed import ImageEmbedding, TextEmbedding

        _image_encoder = ImageEmbedding(model_name=IMAGE_VISION_MODEL)
        _image_text_encoder = TextEmbedding(model_name=IMAGE_TEXT_MODEL)
    return _image_encoder, _image_text_encoder


def _get_prompt_vectors(text_encoder, prompts: list[str]) -> list[list[float]]:
    """Embed CLIP text prompts, caching each one since the prompt set is fixed."""
    missing = [prompt for prompt in dict.fromkeys(prompts) if prompt not in _prompt_vectors]
    if missing:
        for prompt, vector in zip(missing, text_encoder.embed(missing)):
            _prompt_vectors[prompt] = [float(value) for value in vector]
    return [_prompt_vectors[prompt] for prompt in prompts]


def describe_image(data: bytes) -> ImageDescription | None:
    """Embed an image with CLIP and label it against a fixed concept list.

    This is the only source of searchable content for photographs, which produce no
    OCR text at all. The returned vector also powers visual search. Returns None when
    the vision models are unavailable.
    """
    global _last_image_analysis
    import tempfile

    cache_key = sha256_bytes(data)
    if _last_image_analysis and _last_image_analysis[0] == cache_key:
        return _last_image_analysis[1]

    result = _describe_image_uncached(data, tempfile)
    _last_image_analysis = (cache_key, result)
    return result


def _softmax_probabilities(similarities: list[float]) -> list[float]:
    """Turn raw image/text similarities into usable odds via the model's logit scale."""
    scaled = [value * IMAGE_LOGIT_SCALE for value in similarities]
    highest = max(scaled)
    exponentials = [math.exp(value - highest) for value in scaled]
    total = sum(exponentials) or 1.0
    return [value / total for value in exponentials]


def _probe_attributes(
    image_vector: list[float],
    text_encoder,
    group_names: tuple[str, ...],
) -> tuple[list[str], list[str]]:
    """Pick the winning option in each attribute group, skipping ambiguous ones."""
    fragments: list[str] = []
    keywords: list[str] = []

    for group_name in group_names:
        options = _ATTRIBUTE_GROUPS.get(group_name)
        if not options:
            continue

        prompts = [option[0] for option in options]
        vectors = _get_prompt_vectors(text_encoder, prompts)
        similarities = [_cosine_similarity(image_vector, vector) for vector in vectors]
        if not similarities:
            continue

        probabilities = _softmax_probabilities(similarities)
        ranking = sorted(range(len(probabilities)), key=lambda index: probabilities[index], reverse=True)
        best = ranking[0]
        runner_up = probabilities[ranking[1]] if len(ranking) > 1 else 0.0
        if probabilities[best] < ATTRIBUTE_THRESHOLD or probabilities[best] - runner_up < ATTRIBUTE_MARGIN:
            continue

        _, fragment, option_keywords = options[best]
        if fragment:
            fragments.append(fragment)
        keywords.extend(option_keywords)

    return fragments, keywords


def _probe_objects(image_vector: list[float], text_encoder) -> tuple[str, str, list[str]] | None:
    """Name the object in a photo.

    Returns (display name, document type, keywords) or None when no noun stands out.
    """
    prompts = [entry[0] for entry in _OBJECT_NOUNS]
    vectors = _get_prompt_vectors(text_encoder, prompts)
    similarities = [_cosine_similarity(image_vector, vector) for vector in vectors]
    if not similarities:
        return None

    probabilities = _softmax_probabilities(similarities)
    ranking = sorted(range(len(probabilities)), key=lambda index: probabilities[index], reverse=True)
    best = ranking[0]
    best_probability = probabilities[best]
    runner_up_index = ranking[1] if len(ranking) > 1 else None
    runner_up = probabilities[runner_up_index] if runner_up_index is not None else 0.0

    if best_probability < OBJECT_LABEL_THRESHOLD or best_probability - runner_up < OBJECT_LABEL_MARGIN:
        return None

    _, name, document_type, keywords = _OBJECT_NOUNS[best]
    return name, document_type, ["photo", "image", *keywords]


def _compose_caption(base_caption: str, label: str, fragments: list[str]) -> str:
    """Build a caption that distinguishes this image from others of the same kind."""
    if not fragments:
        return base_caption

    if label in {"person_photo", "photo", "event_photo"}:
        # Order the pieces so the result reads naturally, e.g.
        # "Photo of a young man wearing glasses, selfie, indoors".
        subject = next((item for item in fragments if item.startswith(("a ", "an ", "two ", "a group"))), None)
        remaining = [item for item in fragments if item != subject]
        lead = f"Photo of {subject}" if subject else base_caption
        return ", ".join([lead, *remaining])[:200]

    return ", ".join([base_caption, *fragments])[:200]


def _describe_image_uncached(data: bytes, tempfile) -> ImageDescription | None:
    try:
        image_encoder, text_encoder = _get_image_encoders()
        concept_vectors = _get_prompt_vectors(text_encoder, [concept[0] for concept in _IMAGE_CONCEPTS])

        with tempfile.NamedTemporaryFile(suffix=".img", delete=False) as handle:
            handle.write(data)
            temporary_path = handle.name
        try:
            vectors = list(image_encoder.embed([temporary_path]))
        finally:
            try:
                os.unlink(temporary_path)
            except OSError:
                pass

        if not vectors:
            return None
        image_vector = [float(value) for value in vectors[0]]
    except Exception as exc:
        print(f"Local image understanding unavailable: {exc}")
        return None

    similarities = [_cosine_similarity(image_vector, vector) for vector in concept_vectors]
    if not similarities:
        return None

    probabilities = _softmax_probabilities(similarities)
    ranking = sorted(range(len(probabilities)), key=lambda index: probabilities[index], reverse=True)
    best_index = ranking[0]
    best_probability = probabilities[best_index]
    runner_up = probabilities[ranking[1]] if len(ranking) > 1 else 0.0

    ambiguous = (
        best_probability < IMAGE_LABEL_THRESHOLD
        or best_probability - runner_up < IMAGE_LABEL_MARGIN
    )
    if ambiguous:
        # Record it as a photo rather than asserting a wrong label. The vector is still
        # returned so the image remains visually searchable.
        label = "photo"
        document_type = "photo"
        base_caption = "Photograph"
        keywords = ["photo", "image", "picture"]
    else:
        _, label, document_type, base_caption, concept_keywords = _IMAGE_CONCEPTS[best_index]
        keywords = list(concept_keywords)

    # For a generic "thing" photo, name the object. This supersedes the coarse label,
    # which is often simply wrong when the right option was missing from the short list.
    if label in _OBJECT_PROBE_LABELS:
        named = _probe_objects(image_vector, text_encoder)
        if named:
            object_name, object_type, object_keywords = named
            base_caption = f"Photo of a {object_name}"
            document_type = object_type
            label = object_type
            keywords = object_keywords

    # Add distinguishing detail so two images of the same kind do not end up identical.
    fragments, attribute_keywords = _probe_attributes(
        image_vector, text_encoder, _ATTRIBUTE_GROUPS_BY_LABEL.get(label, ())
    )
    caption = _compose_caption(base_caption, label, fragments)
    for keyword in attribute_keywords:
        if keyword not in keywords:
            keywords.append(keyword)

    return ImageDescription(
        label=label,
        document_type=document_type,
        caption=caption,
        keywords=keywords,
        confidence=round(best_probability, 4),
        vector=image_vector,
    )


def embed_image_query(query: str) -> list[float] | None:
    """Embed a text query into CLIP's shared space so it can match image vectors."""
    normalized = re.sub(r"\s+", " ", query or "").strip()
    if not normalized:
        return None
    try:
        _, text_encoder = _get_image_encoders()
        # A short natural-language prompt matches how the concept labels were embedded.
        vectors = list(text_encoder.embed([f"a photo of {normalized}"]))
    except Exception as exc:
        print(f"Visual query embedding unavailable: {exc}")
        return None
    if not vectors:
        return None
    return [float(value) for value in vectors[0]]


def _get_reranker():
    global _reranker
    if _reranker is None:
        from fastembed.rerank.cross_encoder import TextCrossEncoder

        _reranker = TextCrossEncoder(model_name=RERANK_MODEL)
    return _reranker


def rerank_passages(query: str, passages: list[str]) -> list[float] | None:
    """Score query/passage pairs with a cross-encoder.

    Returns calibrated relevance in 0..1, or None when the model cannot be loaded so
    that callers keep the original fusion ordering instead of failing the search.
    """
    cleaned = [re.sub(r"\s+", " ", passage or "").strip()[:2000] or "empty document" for passage in passages]
    if not cleaned:
        return []

    normalized_query = re.sub(r"\s+", " ", query or "").strip()
    if not normalized_query:
        return None

    try:
        encoder = _get_reranker()
        raw_scores = list(encoder.rerank(normalized_query, cleaned))
    except Exception as exc:
        print(f"Local reranker unavailable; keeping fusion order: {exc}")
        return None

    if len(raw_scores) != len(cleaned):
        print("Reranker returned an unexpected score count; keeping fusion order")
        return None

    # Cross-encoder outputs are logits, so squash them into a probability-like range.
    return [round(1.0 / (1.0 + math.exp(-max(-30.0, min(30.0, float(score))))), 4) for score in raw_scores]


def _cosine_similarity(left: list[float] | None, right: list[float] | None) -> float:
    if not left or not right:
        return 0.0
    numerator = sum(a * b for a, b in zip(left, right))
    denominator = math.sqrt(sum(a * a for a in left)) * math.sqrt(sum(b * b for b in right))
    return numerator / denominator if denominator else 0.0


def _folder_candidates(folder_map: dict[str, Any]) -> list[FolderCandidate]:
    candidates: list[FolderCandidate] = []
    for subject, data in (folder_map or {}).items():
        if not isinstance(data, dict) or not data.get("id"):
            continue
        units = data.get("units") or {}
        candidates.append(FolderCandidate(
            folder_id=str(data["id"]),
            label=subject,
            subject=subject,
            unit=None,
        ))
        if isinstance(units, dict):
            for unit, folder_id in units.items():
                if folder_id:
                    candidates.append(FolderCandidate(
                        folder_id=str(folder_id),
                        label=f"{subject} / {unit}",
                        subject=subject,
                        unit=unit,
                    ))
    return candidates


# Maps a detected document signal to folder-name vocabulary. The first tuple holds
# specific names that identify the exact destination; the second holds broader parent
# names that are a reasonable but less certain home.
_SIGNAL_FOLDER_KEYS: dict[str, tuple[tuple[str, ...], tuple[str, ...]]] = {
    "aadhaar": (
        ("aadhaar", "aadhar", "uid", "uidai"),
        ("identity", "identification", "ids", "important", "personal", "government", "kyc"),
    ),
    "pan": (
        ("pan", "permanent"),
        ("identity", "identification", "ids", "important", "personal", "tax", "kyc"),
    ),
    "passport": (
        ("passport",),
        ("identity", "identification", "ids", "important", "travel", "personal", "kyc"),
    ),
    "driving_licence": (
        ("licence", "license", "driving", "dl"),
        ("identity", "identification", "ids", "important", "personal", "vehicle", "kyc"),
    ),
    "voter_id": (
        ("voter", "epic"),
        ("identity", "identification", "ids", "important", "personal", "kyc"),
    ),
    "identity_other": (
        ("identity", "identification", "ids", "kyc"),
        ("important", "personal", "documents"),
    ),
    "screenshot": (
        ("screenshot", "screenshots", "screen", "screens", "captures"),
        ("image", "images", "photo", "photos", "picture", "pictures", "media", "gallery"),
    ),
    "invoice": (
        ("invoice", "invoices", "receipt", "receipts", "bill", "bills"),
        ("finance", "financial", "expense", "expenses", "purchase", "purchases", "payment", "payments", "important"),
    ),
    "bank_statement": (
        ("bank", "banking", "statement", "statements", "passbook"),
        ("finance", "financial", "account", "accounts", "important"),
    ),
    "certificate": (
        ("certificate", "certificates", "certification", "certifications"),
        ("important", "academic", "academics", "education", "award", "awards", "achievements"),
    ),
    "resume": (
        ("resume", "resumes", "cv"),
        ("career", "job", "jobs", "work", "professional", "personal"),
    ),
    "syllabus": (
        ("syllabus", "syllabi", "curriculum"),
        ("academic", "academics", "course", "courses", "education", "study", "studies", "college", "school", "semester"),
    ),
    "assignment": (
        ("assignment", "assignments", "homework"),
        ("academic", "academics", "course", "courses", "education", "study", "studies", "college", "school"),
    ),
    "question_paper": (
        ("question", "questions", "exam", "exams", "paper", "papers", "test", "tests", "pyq"),
        ("academic", "academics", "education", "study", "studies", "college", "school", "semester"),
    ),
    "study_notes": (
        ("notes", "note", "lecture", "lectures"),
        ("academic", "academics", "course", "courses", "education", "study", "studies", "college", "school"),
    ),
    "medical_report": (
        ("medical", "health", "prescription", "prescriptions", "hospital", "reports"),
        ("important", "personal", "insurance"),
    ),
    # Visual-only categories, produced by describe_image for photographs.
    "person_photo": (
        ("photo", "photos", "photographs", "portrait", "portraits", "people", "family", "pictures"),
        ("personal", "image", "images", "gallery", "media", "memories"),
    ),
    "event_photo": (
        ("event", "events", "photos", "photographs", "pictures", "memories", "trip", "trips"),
        ("personal", "image", "images", "gallery", "media"),
    ),
    "scenery_photo": (
        ("scenery", "landscape", "landscapes", "travel", "places", "photos", "photographs", "pictures"),
        ("personal", "image", "images", "gallery", "media"),
    ),
    "food_photo": (
        ("food", "recipes", "meals", "photos", "pictures"),
        ("personal", "image", "images", "gallery", "media"),
    ),
    "vehicle_photo": (
        ("vehicle", "vehicles", "car", "cars", "bike", "bikes"),
        ("personal", "photos", "image", "images", "gallery", "media"),
    ),
    "product_photo": (
        ("product", "products", "shopping", "items"),
        ("personal", "photos", "image", "images", "gallery", "media"),
    ),
    "chart_image": (
        ("chart", "charts", "graphs", "reports", "analytics"),
        ("work", "business", "important", "screenshots"),
    ),
    "source_code": (
        ("code", "coding", "programming", "dev", "development", "snippets", "programs"),
        ("screenshot", "screenshots", "work", "projects", "practice"),
    ),
}

_DOCUMENT_TYPE_TO_SIGNAL = {
    "screenshot": "screenshot",
    "receipt": "invoice",
    "bank_statement": "bank_statement",
    "syllabus": "syllabus",
    "assignment": "assignment",
    "question_paper": "question_paper",
    "study_notes": "study_notes",
    "certificate": "certificate",
    "resume": "resume",
    "medical_report": "medical_report",
    # Visual-only types from describe_image.
    "person_photo": "person_photo",
    "event_photo": "event_photo",
    "scenery_photo": "scenery_photo",
    "food_photo": "food_photo",
    "vehicle_photo": "vehicle_photo",
    "product_photo": "product_photo",
    "chart": "chart_image",
    "document_scan": "study_notes",
    "handwritten_notes": "study_notes",
    "whiteboard_photo": "study_notes",
}


def _label_tokens(text: str) -> set[str]:
    """Tokenise a folder label without stop-word filtering.

    _tokens drops words like "notes" and "unit", which are meaningful folder names,
    and its 3-character minimum would discard "id" and "cv".
    """
    return {token for token in re.findall(r"[a-z0-9]+", (text or "").lower()) if token}


def _document_signals(text: str, filename: str, metadata: DocumentMetadata) -> set[str]:
    """Collect high-precision routing signals for a document."""
    corpus = f"{filename} {metadata.title} {' '.join(metadata.keywords)} {text[:6000]}".lower()
    raw = f"{filename} {text[:6000]}"
    signals: set[str] = set()

    subtype = _identity_subtype(corpus, raw)
    if subtype:
        signals.add(subtype)
        signals.add("identity_other")

    mapped = _DOCUMENT_TYPE_TO_SIGNAL.get(metadata.document_type)
    if mapped:
        signals.add(mapped)
    elif metadata.document_type == "identity_card" and not subtype:
        signals.add("identity_other")

    if _looks_like_source_code(corpus):
        signals.add("source_code")

    return signals


# A broad signal should lose to an exact one when both name a real folder, so that
# "Aadhar Card" beats a generic "Identity Cards" parent for an actual Aadhaar card.
_SIGNAL_SPECIFICITY = {"identity_other": 0.85}

# Words too generic to identify a folder on their own. A unit called "Introduction"
# must not capture every document that happens to contain that word.
_GENERIC_LABEL_WORDS = {
    "and", "basics", "chapter", "concepts", "course", "details", "for", "fundamentals",
    "general", "introduction", "misc", "miscellaneous", "module", "notes", "of", "other",
    "others", "overview", "part", "the", "to", "topics", "unit", "with",
}

# Coverage of a folder's distinctive words that must appear in the document before the
# match counts as deterministic evidence.
NAME_MATCH_STRONG = 0.80
NAME_MATCH_PARTIAL = 0.60

# Suffixes stripped when comparing words, longest first. Without this, a folder named
# "Pipelining" never matches a document about a "pipeline", even though Postgres
# full-text search treats them as the same word.
_STEM_SUFFIXES = (
    "izations", "isations", "ization", "isation", "ational", "tional", "ements",
    "ements", "ations", "ement", "ation", "ments", "ities", "ingly", "edly", "ments",
    "ance", "ence", "ings", "ical", "ally", "ible", "able", "ness", "ment", "ions",
    "ing", "ies", "ied", "ers", "ion", "ity", "ive", "ous", "ful", "est", "als",
    "ed", "es", "er", "ly", "al", "ic", "s", "y",
)


def _stem(token: str) -> str:
    """Reduce a word to a comparable root, mirroring what a text search would do."""
    word = (token or "").lower()
    if len(word) <= 4:
        return word
    for suffix in _STEM_SUFFIXES:
        if word.endswith(suffix) and len(word) - len(suffix) >= 4:
            word = word[: -len(suffix)]
            break
    # "pipeline" and "pipelining" only converge once a trailing vowel is dropped too.
    if len(word) > 4 and word.endswith("e"):
        word = word[:-1]
    return word


def _stems(tokens: Iterable[str]) -> set[str]:
    return {_stem(token) for token in tokens if token}


def _name_match_score(metadata: DocumentMetadata, filename: str, candidate: FolderCandidate) -> float:
    """Score how strongly a document's own title names this folder.

    Syllabus-derived folders are named after chapter titles, and notes usually carry that
    title, so an exact name match is as reliable as document-type detection. Generic
    words are ignored so a folder like "Introduction" cannot match everything.
    """
    from rapidfuzz.fuzz import token_set_ratio

    target = candidate.unit or candidate.subject
    label_tokens = _label_tokens(target) - _GENERIC_LABEL_WORDS
    if not label_tokens:
        return 0.0

    label_stems = _stems(label_tokens)
    # Title and keywords together are the document's prominent vocabulary. Keywords are
    # frequency ranked, so a term appearing there is genuinely about the subject rather
    # than a passing mention.
    prominent_stems = _stems(
        re.findall(r"[a-z0-9]+", (metadata.title or "").lower())
        + [keyword.lower() for keyword in metadata.keywords[:25]]
        + re.findall(r"[a-z0-9]+", Path(filename or "").stem.replace("_", " ").lower())
    )

    # A single-word folder such as "Spark" or "Pipelining" is distinctive enough on its
    # own, provided the word is prominent in the document rather than buried in the body.
    if len(label_stems) < 2:
        only = next(iter(label_stems))
        return 1.0 if len(only) >= 4 and only in prominent_stems else 0.0

    coverage = len(label_stems & prominent_stems) / len(label_stems)
    fuzzy = token_set_ratio(target.lower(), (metadata.title or "").lower()) / 100.0
    best = max(coverage, fuzzy)

    if best >= NAME_MATCH_STRONG:
        return 1.0
    if best >= NAME_MATCH_PARTIAL:
        return 0.6
    return 0.0


# Structural nouns that describe the kind of item rather than qualifying it. An Aadhaar
# card scan rarely contains the word "card", so these must not count as unmet qualifiers.
_CATEGORY_LABEL_WORDS = {
    "card", "cards", "copies", "copy", "doc", "docs", "document", "documents", "file",
    "files", "folder", "folders", "image", "images", "photo", "photos", "picture",
    "pictures", "scan", "scans", "screenshot", "screenshots", "stuff",
}


def _distinctive_label_tokens(label_tokens: set[str]) -> set[str]:
    return label_tokens - _GENERIC_LABEL_WORDS - _CATEGORY_LABEL_WORDS


def _qualifier_support(label_tokens: set[str], matched: set[str], document_tokens: set[str]) -> float:
    """Penalise a folder whose extra qualifier words the document does not support.

    "Coding Screenshot" and "Screenshots" both match a screenshot, but the first also
    claims to be about code. Unless the document mentions that, the broader folder is the
    safer destination.
    """
    extra = _distinctive_label_tokens(label_tokens) - matched
    if not extra:
        return 1.0
    return 1.0 if extra <= document_tokens else 0.75


def _label_specificity(candidate: FolderCandidate, document_tokens: set[str]) -> int:
    """How many of a folder's distinctive words the document actually supports.

    Used to break ties so a precisely named folder wins over a broader one when the
    document backs up the extra words.
    """
    label_tokens = _label_tokens(f"{candidate.subject} {candidate.unit or ''}")
    return len(_distinctive_label_tokens(label_tokens) & document_tokens)


def _rule_folder_score(
    signals: set[str],
    candidate: FolderCandidate,
    document_tokens: set[str] | None = None,
) -> float:
    """Score a folder from deterministic evidence, independent of embeddings."""
    if not signals:
        return 0.0

    subject_tokens = _label_tokens(candidate.subject)
    unit_tokens = _label_tokens(candidate.unit) if candidate.unit else set()
    supporting = document_tokens or set()
    best = 0.0

    for signal in signals:
        specific, generic = _SIGNAL_FOLDER_KEYS.get(signal, ((), ()))
        weight = _SIGNAL_SPECIFICITY.get(signal, 1.0)
        specific_set = set(specific)

        active_tokens = unit_tokens if candidate.unit else subject_tokens
        matched_specific = active_tokens & specific_set
        if matched_specific:
            weight *= _qualifier_support(active_tokens, matched_specific, supporting)

        if candidate.unit:
            # A subfolder must be named for this kind of document itself. It must not
            # inherit credit from a broad parent, otherwise an invoice lands in
            # "Important Documents / PAN Card" purely because the parent is "Important".
            if unit_tokens & set(specific):
                best = max(best, weight)
            elif unit_tokens & set(generic):
                best = max(best, 0.5 * weight)
        else:
            if subject_tokens & set(specific):
                best = max(best, weight)
            elif subject_tokens & set(generic):
                best = max(best, 0.55 * weight)

    return best


def _normalise_by_spread(values: list[float]) -> list[float]:
    """Rescale scores across candidates, or discard them when they carry no signal."""
    if not values:
        return []
    low = min(values)
    high = max(values)
    if high - low < MIN_SEMANTIC_SPREAD:
        return [0.0 for _ in values]
    return [(value - low) / (high - low) for value in values]


def _lexical_folder_score(document_text: str, candidate: FolderCandidate) -> float:
    from rapidfuzz.fuzz import token_set_ratio

    candidate_text = f"{candidate.subject} {candidate.unit or ''}".strip().lower()
    document_tokens = set(_tokens(document_text))
    candidate_tokens = set(_tokens(candidate_text))
    overlap = len(document_tokens & candidate_tokens) / max(len(candidate_tokens), 1)
    fuzzy = token_set_ratio(candidate_text, document_text[:5000].lower()) / 100.0
    return min(1.0, 0.65 * overlap + 0.35 * fuzzy)


def classify_document(
    extracted: ExtractedDocument,
    metadata: DocumentMetadata,
    filename: str,
    folder_map: dict[str, Any],
    root_folder_id: str,
) -> ClassificationResult:
    candidates = _folder_candidates(folder_map)
    fallback = next(
        (candidate for candidate in candidates if candidate.unit is None and candidate.subject.lower() == "imported documents"),
        None,
    )
    # An image understood visually has no text but does have keywords, so it can still
    # be classified. Only give up when there is nothing to reason about at all.
    has_visual_metadata = bool(metadata.keywords) and metadata.document_type not in {"image", "document"}
    if not extracted.text.strip() and not has_visual_metadata:
        target = fallback or FolderCandidate(
            folder_id=root_folder_id,
            label="Imported Documents",
            subject="Imported Documents",
            unit=None,
        )
        return ClassificationResult(
            target_folder_id=target.folder_id,
            target_label=target.label,
            subject=target.subject,
            unit=target.unit,
            confidence=0.0,
            status="unclassified",
        )

    if not candidates:
        return ClassificationResult(
            target_folder_id=root_folder_id,
            target_label="Imported Documents",
            subject="Imported Documents",
            unit=None,
            confidence=0.0,
            status="unclassified",
        )

    # Keep the document descriptor compact so it is comparable in length to a folder
    # descriptor. Embedding 7000 characters against a short label produced near
    # identical similarity for every folder, which made the ranking meaningless.
    document_descriptor = " ".join(filter(None, [
        metadata.title,
        metadata.document_type.replace("_", " "),
        " ".join(metadata.keywords[:12]),
        metadata.summary[:400],
    ])).strip()
    lexical_source = "\n".join([filename, metadata.title, metadata.summary, " ".join(metadata.keywords)])

    candidate_descriptors = [
        f"Subject: {candidate.subject}. Topic or unit: {candidate.unit or candidate.subject}. Folder: {candidate.label}."
        for candidate in candidates
    ]
    vectors = embed_texts([document_descriptor, *candidate_descriptors])
    document_vector = vectors[0] if vectors else None

    signals = _document_signals(extracted.text, filename, metadata)
    semantics = [
        _cosine_similarity(document_vector, vectors[index + 1] if len(vectors) > index + 1 else None)
        for index in range(len(candidates))
    ]
    spread_adjusted = _normalise_by_spread(semantics)

    supporting_tokens = set(re.findall(
        r"[a-z0-9]+",
        f"{filename} {metadata.title} {' '.join(metadata.keywords[:25])} {extracted.text[:4000]}".lower(),
    ))

    if "source_code" in signals:
        # Make the concept searchable by folder names like "Coding Screenshot", since
        # source code rarely contains the words describing it.
        supporting_tokens.update({"code", "coding", "programming"})

    rules: dict[str, float] = {}
    specificities: dict[str, int] = {}
    for index, candidate in enumerate(candidates):
        # Either kind of deterministic evidence counts: the document's type matching a
        # folder's purpose, or the document's title naming the folder.
        rule = max(
            _rule_folder_score(signals, candidate, supporting_tokens),
            _name_match_score(metadata, filename, candidate),
        )
        rules[candidate.folder_id] = rule
        specificities[candidate.folder_id] = _label_specificity(candidate, supporting_tokens)
        lexical = _lexical_folder_score(lexical_source, candidate)
        blended = 0.45 * rule + 0.30 * spread_adjusted[index] + 0.25 * lexical
        candidate.score = round(max(0.0, min(1.0, blended)), 4)

    # Never route into the fallback folder on score alone; it is the destination of
    # last resort, so it competing on merit would suppress real matches.
    contenders = [candidate for candidate in candidates if not fallback or candidate.folder_id != fallback.folder_id]
    # Rank by rule evidence first and similarity only as a tie-break, so a folder named
    # for exactly this document type cannot be displaced by a vaguer folder that merely
    # scores slightly higher on embeddings.
    ranked = sorted(
        contenders,
        key=lambda item: (
            rules.get(item.folder_id, 0.0),
            specificities.get(item.folder_id, 0),
            item.score,
        ),
        reverse=True,
    )

    target = fallback or FolderCandidate(
        folder_id=root_folder_id,
        label="Imported Documents",
        subject="Imported Documents",
        unit=None,
    )

    if not ranked:
        return ClassificationResult(
            target_folder_id=target.folder_id,
            target_label=target.label,
            subject=target.subject,
            unit=target.unit,
            confidence=0.0,
            status="automatic",
        )

    best = ranked[0]
    second = ranked[1] if len(ranked) > 1 else None
    second_score = second.score if second else 0.0

    # An unambiguous deterministic match is trusted directly; a single folder named
    # for exactly this document type is a stronger signal than any similarity score.
    exclusive_rule_match = rules.get(best.folder_id, 0.0) >= 1.0 and (
        second is None
        or rules.get(second.folder_id, 0.0) < 1.0
        # Equally strong rules are still decisive when one folder's name is more
        # precisely supported by the document.
        or specificities.get(best.folder_id, 0) > specificities.get(second.folder_id, 0)
    )
    confident = exclusive_rule_match or (
        best.score >= AUTO_CLASSIFICATION_THRESHOLD
        and best.score - second_score >= AUTO_CLASSIFICATION_MARGIN
    )

    def worth_offering(candidate: FolderCandidate) -> bool:
        if candidate.score < ALTERNATIVE_SCORE_FLOOR:
            return False
        # Similarity alone is weak evidence, so a folder with no rule support has to
        # clear a higher bar before it is presented as a real option.
        if rules.get(candidate.folder_id, 0.0) > 0.0:
            return True
        return candidate.score >= STRONG_SIMILARITY_SCORE

    if confident:
        return ClassificationResult(
            target_folder_id=best.folder_id,
            target_label=best.label,
            subject=best.subject,
            unit=best.unit,
            confidence=best.score,
            status="automatic",
            alternatives=[item for item in ranked[1:3] if worth_offering(item)],
        )

    # Only offer choices backed by evidence. Previously the top three folders were
    # always shown, so an unrelated file was offered whatever happened to rank first.
    alternatives = [candidate for candidate in ranked if worth_offering(candidate)][:3]
    return ClassificationResult(
        target_folder_id=target.folder_id,
        target_label=target.label,
        subject=target.subject,
        unit=target.unit,
        confidence=best.score,
        status="needs_confirmation" if alternatives else "automatic",
        alternatives=alternatives,
    )
